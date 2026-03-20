"""Knowledge Base & RAG API routes."""

import uuid
import hashlib
import logging
import os
from datetime import datetime, timezone

from fastapi import APIRouter, Depends, HTTPException, Query, UploadFile, File
from pydantic import BaseModel, Field
from sqlalchemy.ext.asyncio import AsyncSession

from app.database import get_db
from app.api.deps import get_current_user_id
from app.services.rag_service import rag_service, kb_service
from app.models.knowledge_base import KnowledgeBase, KnowledgeDocument

logger = logging.getLogger(__name__)

router = APIRouter(prefix="/admin/knowledge-bases", tags=["Admin - Knowledge Base"])

ALLOWED_DOC_EXTENSIONS = {".pdf", ".docx", ".doc", ".xlsx", ".xls", ".pptx", ".ppt", ".txt", ".md", ".csv", ".html", ".json", ".xml"}


# ─── Schemas ───

class KBCreate(BaseModel):
    name: str = Field(..., min_length=1, max_length=200)
    description: str | None = None
    department: str | None = None
    is_public: bool = False
    allowed_roles: list | None = None
    embedding_model: str = Field(default="nomic-embed-text")
    chunk_size: int = Field(default=500, ge=100, le=5000)
    chunk_overlap: int = Field(default=50, ge=0, le=500)


class KBUpdate(BaseModel):
    name: str | None = None
    description: str | None = None
    department: str | None = None
    is_public: bool | None = None
    allowed_roles: list | None = None
    chunk_size: int | None = None
    chunk_overlap: int | None = None


def _serialize_kb(kb) -> dict:
    return {
        "id": str(kb.id),
        "name": kb.name,
        "description": kb.description,
        "department": kb.department,
        "is_public": kb.is_public,
        "allowed_roles": kb.allowed_roles,
        "embedding_model": kb.embedding_model,
        "chunk_size": kb.chunk_size,
        "chunk_overlap": kb.chunk_overlap,
        "document_count": kb.document_count,
        "total_chunks": kb.total_chunks,
        "last_synced_at": kb.last_synced_at.isoformat() if kb.last_synced_at else None,
        "created_at": kb.created_at.isoformat() if kb.created_at else None,
    }


def _serialize_doc(doc) -> dict:
    return {
        "id": str(doc.id),
        "knowledge_base_id": str(doc.knowledge_base_id),
        "title": doc.title,
        "file_name": doc.file_name,
        "file_type": doc.file_type,
        "file_size": doc.file_size,
        "status": doc.status,
        "chunk_count": doc.chunk_count,
        "error_message": doc.error_message,
        "version": doc.version,
        "created_at": doc.created_at.isoformat() if doc.created_at else None,
    }


# ─── Routes ───

@router.get("")
async def list_knowledge_bases(db: AsyncSession = Depends(get_db)):
    kbs = await kb_service.list_knowledge_bases(db)
    return {"knowledge_bases": [_serialize_kb(kb) for kb in kbs]}


@router.post("")
async def create_knowledge_base(
    body: KBCreate,
    user_id: uuid.UUID = Depends(get_current_user_id),
    db: AsyncSession = Depends(get_db),
):
    data = body.model_dump(exclude_none=True)
    data["created_by"] = user_id
    kb = await kb_service.create_knowledge_base(data, db)
    return _serialize_kb(kb)


@router.get("/stats")
async def kb_stats(db: AsyncSession = Depends(get_db)):
    return await kb_service.get_kb_stats(db)


@router.get("/{kb_id}")
async def get_knowledge_base(kb_id: uuid.UUID, db: AsyncSession = Depends(get_db)):
    kb = await kb_service.get_knowledge_base(kb_id, db)
    if not kb:
        raise HTTPException(status_code=404, detail="Knowledge base not found")
    return _serialize_kb(kb)


@router.patch("/{kb_id}")
async def update_knowledge_base(
    kb_id: uuid.UUID,
    body: KBUpdate,
    db: AsyncSession = Depends(get_db),
):
    data = body.model_dump(exclude_none=True)
    kb = await kb_service.update_knowledge_base(kb_id, data, db)
    if not kb:
        raise HTTPException(status_code=404, detail="Knowledge base not found")
    return _serialize_kb(kb)


@router.delete("/{kb_id}")
async def delete_knowledge_base(kb_id: uuid.UUID, db: AsyncSession = Depends(get_db)):
    if not await kb_service.delete_knowledge_base(kb_id, db):
        raise HTTPException(status_code=404, detail="Knowledge base not found")
    return {"status": "deleted"}


@router.post("/{kb_id}/documents")
async def upload_document(
    kb_id: uuid.UUID,
    file: UploadFile = File(...),
    user_id: uuid.UUID = Depends(get_current_user_id),
    db: AsyncSession = Depends(get_db),
):
    """Upload and process a document into the knowledge base."""
    kb = await kb_service.get_knowledge_base(kb_id, db)
    if not kb:
        raise HTTPException(status_code=404, detail="Knowledge base not found")

    # Validate extension
    ext = os.path.splitext(file.filename or "")[1].lower()
    if ext not in ALLOWED_DOC_EXTENSIONS:
        raise HTTPException(status_code=400, detail=f"File type {ext} not supported")

    # Read content
    content_bytes = await file.read()
    file_hash = hashlib.sha256(content_bytes).hexdigest()

    # Extract text
    text = await _extract_text(content_bytes, ext)
    if not text:
        raise HTTPException(status_code=400, detail="Could not extract text from file")

    # Create document record
    doc = KnowledgeDocument(
        knowledge_base_id=kb_id,
        title=os.path.splitext(file.filename or "document")[0],
        file_name=file.filename or "unknown",
        file_type=ext.lstrip("."),
        file_size=len(content_bytes),
        file_hash=file_hash,
        uploaded_by=user_id,
    )
    db.add(doc)
    await db.flush()

    # Ingest: chunk + embed
    try:
        await rag_service.ingest_document(doc, text, db)
    except Exception as e:
        doc.status = "failed"
        doc.error_message = str(e)[:500]
        logger.error("Document ingestion failed: %s", str(e))

    return _serialize_doc(doc)


@router.delete("/{kb_id}/documents/{doc_id}")
async def delete_document(
    kb_id: uuid.UUID,
    doc_id: uuid.UUID,
    db: AsyncSession = Depends(get_db),
):
    from sqlalchemy import select, func
    result = await db.execute(
        select(KnowledgeDocument).where(
            KnowledgeDocument.id == doc_id,
            KnowledgeDocument.knowledge_base_id == kb_id,
        )
    )
    doc = result.scalar_one_or_none()
    if not doc:
        raise HTTPException(status_code=404, detail="Document not found")
    await db.delete(doc)
    await db.flush()

    # Update KB counters
    kb = await db.get(KnowledgeBase, kb_id)
    if kb:
        from app.models.knowledge_base import DocumentChunk
        total_docs = (await db.execute(
            select(func.count()).select_from(KnowledgeDocument)
            .where(KnowledgeDocument.knowledge_base_id == kb_id)
        )).scalar() or 0
        total_chunks = (await db.execute(
            select(func.count()).select_from(DocumentChunk)
            .join(KnowledgeDocument)
            .where(KnowledgeDocument.knowledge_base_id == kb_id)
        )).scalar() or 0
        kb.document_count = total_docs
        kb.total_chunks = total_chunks
        await db.flush()

    return {"status": "deleted"}


@router.post("/{kb_id}/sync")
async def sync_knowledge_base(
    kb_id: uuid.UUID,
    db: AsyncSession = Depends(get_db),
):
    """Re-process all documents in the knowledge base."""
    from sqlalchemy import select
    kb = await kb_service.get_knowledge_base(kb_id, db)
    if not kb:
        raise HTTPException(status_code=404, detail="Knowledge base not found")

    result = await db.execute(
        select(KnowledgeDocument).where(KnowledgeDocument.knowledge_base_id == kb_id)
    )
    docs = result.scalars().all()
    errors = []
    for doc in docs:
        try:
            # Re-extract text would require stored content; re-embed existing chunks
            from app.models.knowledge_base import DocumentChunk
            chunk_result = await db.execute(
                select(DocumentChunk).where(DocumentChunk.document_id == doc.id)
                .order_by(DocumentChunk.chunk_index)
            )
            chunks = chunk_result.scalars().all()
            for chunk in chunks:
                embedding = await rag_service.embed_text(chunk.content)
                chunk.embedding = embedding if embedding else None
            doc.status = "ready"
        except Exception as e:
            doc.status = "failed"
            doc.error_message = str(e)[:500]
            errors.append(str(doc.id))

    kb.last_synced_at = datetime.now(timezone.utc)
    await db.flush()
    return {"status": "sync_complete", "document_count": len(docs), "errors": errors}


@router.get("/{kb_id}/search")
async def search_knowledge_base(
    kb_id: uuid.UUID,
    q: str = Query(..., min_length=1),
    top_k: int = Query(5, ge=1, le=20),
    db: AsyncSession = Depends(get_db),
):
    results = await rag_service.search(kb_id, q, db, top_k=top_k)
    return {"results": results, "query": q}


@router.get("/{kb_id}/documents")
async def list_documents(
    kb_id: uuid.UUID,
    db: AsyncSession = Depends(get_db),
):
    from sqlalchemy import select
    result = await db.execute(
        select(KnowledgeDocument)
        .where(KnowledgeDocument.knowledge_base_id == kb_id)
        .order_by(KnowledgeDocument.created_at.desc())
    )
    docs = result.scalars().all()
    return {"documents": [_serialize_doc(d) for d in docs]}


# ─── Text Extraction ───

async def _extract_text(content: bytes, ext: str) -> str:
    """Extract text from various document formats (runs in thread to avoid blocking)."""
    import asyncio
    return await asyncio.to_thread(_extract_text_sync, content, ext)


def _extract_text_sync(content: bytes, ext: str) -> str:
    """Synchronous text extraction from various document formats."""
    import io

    try:
        if ext in (".txt", ".md", ".csv"):
            return content.decode("utf-8", errors="replace")

        if ext == ".pdf":
            try:
                import pdfplumber
                with pdfplumber.open(io.BytesIO(content)) as pdf:
                    return "\n".join(page.extract_text() or "" for page in pdf.pages)
            except ImportError:
                from PyPDF2 import PdfReader
                reader = PdfReader(io.BytesIO(content))
                return "\n".join(page.extract_text() or "" for page in reader.pages)

        if ext in (".docx", ".doc"):
            from docx import Document
            doc = Document(io.BytesIO(content))
            return "\n".join(p.text for p in doc.paragraphs)

        if ext in (".xlsx", ".xls"):
            from openpyxl import load_workbook
            wb = load_workbook(io.BytesIO(content), read_only=True, data_only=True)
            lines = []
            for ws in wb.worksheets:
                for row in ws.iter_rows(values_only=True):
                    lines.append(" | ".join(str(c) if c is not None else "" for c in row))
            return "\n".join(lines)

        if ext in (".pptx", ".ppt"):
            from pptx import Presentation
            prs = Presentation(io.BytesIO(content))
            text_parts = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text_parts.append(shape.text)
            return "\n".join(text_parts)

        if ext in (".html",):
            return content.decode("utf-8", errors="replace")

        if ext in (".json", ".xml"):
            return content.decode("utf-8", errors="replace")

        return content.decode("utf-8", errors="replace")

    except Exception as e:
        logger.error("Text extraction failed for %s: %s", ext, str(e))
        return ""
