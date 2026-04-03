"""Chat / messaging API routes with streaming support."""

import asyncio
import base64
import uuid
import io
import csv
import logging
import re
from html import escape

from fastapi import APIRouter, Depends, Form, HTTPException, Request, UploadFile, File
from fastapi.responses import Response, StreamingResponse
from sqlalchemy import select
from sqlalchemy.ext.asyncio import AsyncSession

from app.database import get_db
from app.api.deps import get_current_user_id, get_current_user_token, get_client_ip
from app.schemas.message import (
    ChatRequest,
    ChatResponse,
    GenerateFileResponse,
    GenerateFileRequest,
    MessageResponse,
    SearchRequest,
    SearchResult,
)
from app.services.chat_service import chat_service
from app.services.audit_service import audit_service
from app.config import settings as app_settings
from app.models.file_upload import FileUpload
from app.models.generated_file import GeneratedFile
from app.models.message import Message

logger = logging.getLogger(__name__)

router = APIRouter(prefix="/chat", tags=["Chat"])

ALLOWED_EXTENSIONS = {
    ".txt", ".csv", ".md",
    ".pdf",
    ".docx", ".doc",
    ".xlsx", ".xls",
    ".pptx", ".ppt",
    ".rtf",
    ".json", ".xml", ".html",
    # Images
    ".jpg", ".jpeg", ".png", ".gif", ".webp", ".bmp", ".svg",
}

IMAGE_EXTENSIONS = {".jpg", ".jpeg", ".png", ".gif", ".webp", ".bmp", ".svg"}

MIME_MAP = {
    ".jpg": "image/jpeg",
    ".jpeg": "image/jpeg",
    ".png": "image/png",
    ".gif": "image/gif",
    ".webp": "image/webp",
    ".bmp": "image/bmp",
    ".svg": "image/svg+xml",
}


@router.post("", response_model=ChatResponse)
async def send_message(
    body: ChatRequest,
    request: Request,
    user_id: uuid.UUID = Depends(get_current_user_id),
    token=Depends(get_current_user_token),
    db: AsyncSession = Depends(get_db),
):
    """Send a message and get a non-streaming LLM response."""
    try:
        conv_id = uuid.UUID(body.conversation_id) if body.conversation_id else None
    except ValueError:
        raise HTTPException(status_code=400, detail="Invalid conversation_id format")

    try:
        agent_id = uuid.UUID(body.agent_id) if body.agent_id else None
    except ValueError:
        raise HTTPException(status_code=400, detail="Invalid agent_id format")

    try:
        user_msg, assistant_msg, conv = await chat_service.send_message(
            user_id=user_id,
            conversation_id=conv_id,
            content=body.message,
            model=body.model,
            db=db,
            agent_id=agent_id,
            deep_analysis=body.deep_analysis,
            vision_images=body.vision_images,
        )
    except PermissionError:
        raise HTTPException(status_code=404, detail="Conversation not found")
    except Exception:
        raise HTTPException(status_code=500, detail="Failed to generate response")

    await audit_service.log(
        db,
        action="prompt_submitted",
        user_id=user_id,
        username=token.username,
        resource_type="conversation",
        resource_id=str(conv.id),
        ip_address=get_client_ip(request),
    )

    return ChatResponse(
        message=MessageResponse(
            id=str(assistant_msg.id),
            conversation_id=str(conv.id),
            role=assistant_msg.role,
            content=assistant_msg.content,
            model=assistant_msg.model,
            token_count=assistant_msg.token_count,
            citations=getattr(assistant_msg, "_citations", []),
            quality_issues=getattr(assistant_msg, "_quality_issues", []),
            followups=getattr(assistant_msg, "_followups", []),
            created_at=assistant_msg.created_at,
        ),
        conversation_id=str(conv.id),
    )


@router.post("/stream")
async def send_message_stream(
    body: ChatRequest,
    request: Request,
    user_id: uuid.UUID = Depends(get_current_user_id),
    token=Depends(get_current_user_token),
    db: AsyncSession = Depends(get_db),
):
    """Send a message and stream the LLM response (SSE-like NDJSON)."""
    try:
        conv_id = uuid.UUID(body.conversation_id) if body.conversation_id else None
    except ValueError:
        raise HTTPException(status_code=400, detail="Invalid conversation_id format")

    try:
        agent_id = uuid.UUID(body.agent_id) if body.agent_id else None
    except ValueError:
        raise HTTPException(status_code=400, detail="Invalid agent_id format")

    await audit_service.log(
        db,
        action="prompt_submitted",
        user_id=user_id,
        username=token.username,
        resource_type="conversation",
        resource_id=str(conv_id) if conv_id else "new",
        ip_address=get_client_ip(request),
    )

    return StreamingResponse(
        chat_service.send_message_stream(
            user_id=user_id,
            conversation_id=conv_id,
            content=body.message,
            model=body.model,
            agent_id=agent_id,
            deep_analysis=body.deep_analysis,
            vision_images=body.vision_images,
        ),
        media_type="application/x-ndjson",
        headers={"Cache-Control": "no-cache", "X-Accel-Buffering": "no"},
    )


@router.post("/search", response_model=list[SearchResult])
async def search_messages(
    body: SearchRequest,
    user_id: uuid.UUID = Depends(get_current_user_id),
    db: AsyncSession = Depends(get_db),
):
    """Search across the current user's messages only."""
    results = await chat_service.search_messages(user_id, body.query, body.limit, db)
    return [SearchResult(**r) for r in results]


@router.get("/models")
async def list_available_models(
    _user_id: uuid.UUID = Depends(get_current_user_id),
):
    """List available LLM models for the model selector (any authenticated user)."""
    from app.services.llm_service import llm_service
    models = await llm_service.list_models()
    return {
        "models": [m.get("name", "unknown") for m in models],
        "default": llm_service.default_model,
    }


def _get_extension(filename: str) -> str:
    """Return lowercase file extension including the dot."""
    import os
    return os.path.splitext(filename)[1].lower()


def _extract_text_txt(data: bytes) -> str:
    return data.decode("utf-8", errors="replace")


def _extract_text_csv(data: bytes) -> str:
    text = data.decode("utf-8", errors="replace")
    try:
        reader = csv.reader(io.StringIO(text))
        rows = []
        for row in reader:
            rows.append(" | ".join(row))
        return "\n".join(rows)
    except Exception:
        # Malformed CSV — return as plain text
        return text


def _extract_text_pdf(data: bytes) -> str:
    import PyPDF2
    try:
        reader = PyPDF2.PdfReader(io.BytesIO(data))
        pages = []
        for page in reader.pages:
            try:
                t = page.extract_text()
                if t:
                    pages.append(t)
            except Exception:
                continue  # skip unreadable pages
        if pages:
            return "\n\n".join(pages)
    except Exception:
        pass
    # Fallback: brute-force extract readable strings from raw PDF bytes
    import re as _re
    text_chunks = _re.findall(rb'[\x20-\x7E]{4,}', data)
    fallback = " ".join(chunk.decode("ascii", errors="ignore") for chunk in text_chunks)
    if fallback.strip():
        return fallback
    raise ValueError("Could not extract any text from PDF")


def _extract_text_docx(data: bytes) -> str:
    import docx
    try:
        doc = docx.Document(io.BytesIO(data))
        text = "\n".join(p.text for p in doc.paragraphs if p.text.strip())
        if text.strip():
            return text
    except Exception:
        pass
    # Fallback for corrupted .docx or old .doc — extract readable strings
    import re as _re
    text_chunks = _re.findall(rb'[\x20-\x7E]{4,}', data)
    fallback = " ".join(chunk.decode("ascii", errors="ignore") for chunk in text_chunks)
    if fallback.strip():
        return fallback
    raise ValueError("Could not extract any text from document")


def _extract_text_xlsx(data: bytes) -> str:
    import openpyxl
    try:
        wb = openpyxl.load_workbook(io.BytesIO(data), read_only=True, data_only=True)
        lines = []
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            lines.append(f"--- Sheet: {sheet_name} ---")
            for row in ws.iter_rows(values_only=True):
                cells = [str(c) if c is not None else "" for c in row]
                lines.append(" | ".join(cells))
        wb.close()
        if lines:
            return "\n".join(lines)
    except Exception:
        pass
    # Fallback: extract readable strings from raw bytes
    import re as _re
    text_chunks = _re.findall(rb'[\x20-\x7E]{4,}', data)
    fallback = " ".join(chunk.decode("ascii", errors="ignore") for chunk in text_chunks)
    if fallback.strip():
        return fallback
    raise ValueError("Could not extract any text from spreadsheet")


def _extract_text_pptx(data: bytes) -> str:
    from pptx import Presentation
    try:
        prs = Presentation(io.BytesIO(data))
        slides_text = []
        for i, slide in enumerate(prs.slides, 1):
            parts = [f"--- Slide {i} ---"]
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        text = paragraph.text.strip()
                        if text:
                            parts.append(text)
            slides_text.append("\n".join(parts))
        if slides_text:
            return "\n\n".join(slides_text)
    except Exception:
        pass
    # Fallback: extract readable strings from raw bytes
    import re as _re
    text_chunks = _re.findall(rb'[\x20-\x7E]{4,}', data)
    fallback = " ".join(chunk.decode("ascii", errors="ignore") for chunk in text_chunks)
    if fallback.strip():
        return fallback
    raise ValueError("Could not extract any text from presentation")


EXTRACTOR_MAP = {
    ".txt": _extract_text_txt,
    ".md": _extract_text_txt,
    ".json": _extract_text_txt,
    ".xml": _extract_text_txt,
    ".html": _extract_text_txt,
    ".rtf": _extract_text_txt,
    ".csv": _extract_text_csv,
    ".pdf": _extract_text_pdf,
    ".docx": _extract_text_docx,
    ".doc": _extract_text_docx,
    ".xlsx": _extract_text_xlsx,
    ".xls": _extract_text_xlsx,
    ".pptx": _extract_text_pptx,
    ".ppt": _extract_text_pptx,
}


@router.get("/attachments-enabled")
async def check_attachments_enabled(
    _user_id: uuid.UUID = Depends(get_current_user_id),
):
    """Check if file attachments are enabled by admin."""
    return {
        "enabled": app_settings.ATTACHMENTS_ENABLED,
        "max_size_mb": app_settings.ATTACHMENTS_MAX_SIZE_MB,
    }


def _sanitize_filename(name: str) -> str:
    """Sanitize a filename for safe display — allow only safe characters."""
    # Keep only alphanumeric, dots, dashes, underscores, spaces
    sanitized = re.sub(r"[^\w.\- ]", "_", name)
    # Collapse consecutive dots (prevent path traversal like ..)
    sanitized = re.sub(r"\.{2,}", ".", sanitized)
    return sanitized.strip() or "unnamed"


def _normalize_output_format(fmt: str) -> str:
    normalized = (fmt or "").strip().lower()
    if normalized == "doc":
        return "docx"
    if normalized == "excel":
        return "xlsx"
    return normalized


def _markdown_to_text(content: str) -> str:
    text = content or ""
    text = re.sub(r"```([\s\S]*?)```", lambda m: m.group(1).strip(), text)
    text = re.sub(r"`([^`]*)`", r"\1", text)
    text = re.sub(r"!\[[^\]]*\]\([^)]*\)", "", text)
    text = re.sub(r"\[([^\]]+)\]\([^)]*\)", r"\1", text)
    text = re.sub(r"^\s{0,3}#{1,6}\s*", "", text, flags=re.MULTILINE)
    text = re.sub(r"^\s*>\s?", "", text, flags=re.MULTILINE)
    text = re.sub(r"^\s*[-*+]\s+", "- ", text, flags=re.MULTILINE)
    text = re.sub(r"^\s*\d+\.\s+", "- ", text, flags=re.MULTILINE)
    return text.strip()


def _build_html_document(content: str) -> str:
    return (
        "<!doctype html>"
        "<html><head><meta charset=\"utf-8\"/>"
        "<title>Generated Response</title>"
        "<style>body{font-family:Segoe UI,Arial,sans-serif;padding:24px;line-height:1.5;}"
        "pre{white-space:pre-wrap;word-break:break-word;background:#f7f7f7;padding:16px;border-radius:8px;}"
        "</style></head><body>"
        "<h1>Generated Response</h1>"
        f"<pre>{escape(content or '')}</pre>"
        "</body></html>"
    )


def _build_generated_file(content: str, output_format: str) -> tuple[bytes, str, str]:
    plain = _markdown_to_text(content)

    if output_format == "pdf":
        from fpdf import FPDF

        pdf = FPDF()
        pdf.set_auto_page_break(auto=True, margin=15)
        pdf.add_page()
        pdf.set_font("Helvetica", "", 11)
        safe_text = plain.encode("latin-1", errors="replace").decode("latin-1")
        pdf.multi_cell(0, 6, safe_text or " ")
        return bytes(pdf.output()), "application/pdf", "pdf"

    if output_format == "docx":
        from docx import Document

        doc = Document()
        for block in (plain or "").split("\n\n"):
            doc.add_paragraph(block.strip() or " ")
        buf = io.BytesIO()
        doc.save(buf)
        return (
            buf.getvalue(),
            "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
            "docx",
        )

    if output_format == "xlsx":
        from openpyxl import Workbook

        wb = Workbook()
        ws = wb.active
        ws.title = "Response"
        lines = plain.splitlines() or [plain or ""]
        for idx, line in enumerate(lines, start=1):
            ws.cell(row=idx, column=1, value=line)
        buf = io.BytesIO()
        wb.save(buf)
        return (
            buf.getvalue(),
            "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
            "xlsx",
        )

    if output_format == "html":
        html = _build_html_document(content)
        return html.encode("utf-8"), "text/html; charset=utf-8", "html"

    raise ValueError("Unsupported output format")


@router.post("/generate-file")
async def generate_file(
    body: GenerateFileRequest,
    user_id: uuid.UUID = Depends(get_current_user_id),
    db: AsyncSession = Depends(get_db),
):
    output_format = _normalize_output_format(body.format)
    try:
        conversation_id = uuid.UUID(body.conversation_id)
        message_id = uuid.UUID(body.message_id)
    except ValueError as exc:
        raise HTTPException(status_code=400, detail="Invalid conversation_id or message_id") from exc

    msg = (
        await db.execute(
            select(Message)
            .where(
                Message.id == message_id,
                Message.conversation_id == conversation_id,
            )
        )
    ).scalar_one_or_none()
    if not msg:
        raise HTTPException(status_code=404, detail="Message not found")

    conv = await chat_service.get_conversation(conversation_id, user_id, db)
    if not conv:
        raise HTTPException(status_code=404, detail="Conversation not found")

    try:
        file_bytes, media_type, extension = _build_generated_file(body.content, output_format)
    except ValueError as exc:
        raise HTTPException(status_code=400, detail=str(exc)) from exc
    except Exception:
        raise HTTPException(status_code=500, detail="Failed to generate file")

    raw_name = body.filename or f"assistant-output.{extension}"
    safe_name = _sanitize_filename(raw_name)
    if not safe_name.lower().endswith(f".{extension}"):
        safe_name = f"{safe_name}.{extension}"

    generated = GeneratedFile(
        user_id=user_id,
        conversation_id=conversation_id,
        message_id=message_id,
        filename=safe_name,
        extension=f".{extension}",
        mime_type=media_type.split(";")[0],
        size_bytes=len(file_bytes),
        content=file_bytes,
    )
    db.add(generated)
    await db.flush()

    return GenerateFileResponse(
        id=str(generated.id),
        name=safe_name,
        type="document",
        url=f"/api/chat/generated-files/{generated.id}/download",
    )


@router.get("/generated-files/{file_id}/download")
async def download_generated_file(
    file_id: uuid.UUID,
    user_id: uuid.UUID = Depends(get_current_user_id),
    db: AsyncSession = Depends(get_db),
):
    file_record = (
        await db.execute(
            select(GeneratedFile).where(
                GeneratedFile.id == file_id,
                GeneratedFile.user_id == user_id,
            )
        )
    ).scalar_one_or_none()
    if not file_record:
        raise HTTPException(status_code=404, detail="Generated file not found")

    return Response(
        content=file_record.content,
        media_type=file_record.mime_type,
        headers={"Content-Disposition": f'attachment; filename="{file_record.filename}"'},
    )


@router.post("/upload")
async def upload_file(
    file: UploadFile = File(...),
    conversation_id: str | None = Form(None),
    user_id: uuid.UUID = Depends(get_current_user_id),
    db: AsyncSession = Depends(get_db),
):
    """Upload a document and extract its text content."""
    logger.info("Upload request received: filename=%s, user=%s", getattr(file, 'filename', 'N/A'), user_id)
    if not app_settings.ATTACHMENTS_ENABLED:
        raise HTTPException(status_code=403, detail="File attachments are disabled by administrator")

    if not file.filename:
        raise HTTPException(status_code=400, detail="No filename provided")

    ext = _get_extension(file.filename)
    if ext not in ALLOWED_EXTENSIONS:
        raise HTTPException(
            status_code=400,
            detail=f"Unsupported file type '{ext}'. Allowed: {', '.join(sorted(ALLOWED_EXTENSIONS))}",
        )

    data = await file.read()
    max_bytes = app_settings.ATTACHMENTS_MAX_SIZE_MB * 1024 * 1024
    if len(data) > max_bytes:
        raise HTTPException(status_code=400, detail=f"File too large. Maximum size is {app_settings.ATTACHMENTS_MAX_SIZE_MB} MB.")

    # --- Image uploads: return base64 data URL instead of extracted text ---
    if ext in IMAGE_EXTENSIONS:
        mime = MIME_MAP.get(ext, "application/octet-stream")
        b64 = base64.b64encode(data).decode("ascii")
        image_url = f"data:{mime};base64,{b64}"

        # Log to DB
        conv_uuid = None
        if conversation_id:
            try:
                conv_uuid = uuid.UUID(conversation_id)
            except ValueError:
                pass
        try:
            upload_record = FileUpload(
                user_id=user_id,
                conversation_id=conv_uuid,
                filename=file.filename,
                extension=ext,
                size_bytes=len(data),
                char_count=0,
                truncated=False,
            )
            db.add(upload_record)
            await db.flush()
        except Exception:
            logger.warning("Failed to log upload to DB for %s (table may not exist)", file.filename)
            await db.rollback()

        return {
            "filename": _sanitize_filename(file.filename),
            "size": len(data),
            "extension": ext,
            "text": "",
            "truncated": False,
            "image_url": image_url,
        }

    # --- Document uploads: extract text ---
    extractor = EXTRACTOR_MAP.get(ext)
    if not extractor:
        raise HTTPException(status_code=400, detail=f"No text extractor for '{ext}'")

    try:
        # Run synchronous extraction in a thread to avoid blocking the event loop (P4)
        loop = asyncio.get_running_loop()
        extracted_text = await loop.run_in_executor(None, extractor, data)
    except Exception:
        logger.exception("Text extraction failed for file: %s", file.filename)
        raise HTTPException(status_code=422, detail="Failed to extract text from the uploaded file")

    # Truncate very long documents to avoid overloading the LLM context
    max_chars = app_settings.ATTACHMENTS_MAX_EXTRACT_CHARS
    truncated = len(extracted_text) > max_chars
    if truncated:
        extracted_text = extracted_text[:max_chars]

    # Log the upload to DB (F1: set conversation_id when available)
    conv_uuid = None
    if conversation_id:
        try:
            conv_uuid = uuid.UUID(conversation_id)
        except ValueError:
            pass
    try:
        upload_record = FileUpload(
            user_id=user_id,
            conversation_id=conv_uuid,
            filename=file.filename,
            extension=ext,
            size_bytes=len(data),
            char_count=len(extracted_text),
            truncated=truncated,
        )
        db.add(upload_record)
        await db.flush()
    except Exception:
        logger.warning("Failed to log upload to DB for %s (table may not exist)", file.filename)
        await db.rollback()

    return {
        "filename": _sanitize_filename(file.filename),
        "size": len(data),
        "extension": ext,
        "text": extracted_text,
        "truncated": truncated,
    }


@router.post("/upload-multiple")
async def upload_multiple_files(
    files: list[UploadFile] = File(...),
    conversation_id: str | None = Form(None),
    user_id: uuid.UUID = Depends(get_current_user_id),
    db: AsyncSession = Depends(get_db),
):
    """Upload multiple documents and extract text from each."""
    if not app_settings.ATTACHMENTS_ENABLED:
        raise HTTPException(status_code=403, detail="File attachments are disabled by administrator")

    if len(files) > 10:
        raise HTTPException(status_code=400, detail="Maximum 10 files per upload")

    results = []
    loop = asyncio.get_running_loop()

    for file in files:
        if not file.filename:
            results.append({"filename": "unknown", "error": "No filename provided"})
            continue

        safe_name = _sanitize_filename(file.filename)
        ext = _get_extension(file.filename)
        if ext not in ALLOWED_EXTENSIONS:
            results.append({"filename": safe_name, "error": f"Unsupported file type '{ext}'"})
            continue

        data = await file.read()
        max_bytes = app_settings.ATTACHMENTS_MAX_SIZE_MB * 1024 * 1024
        if len(data) > max_bytes:
            results.append({"filename": safe_name, "error": f"File too large (max {app_settings.ATTACHMENTS_MAX_SIZE_MB} MB)"})
            continue

        extractor = EXTRACTOR_MAP.get(ext)
        if not extractor:
            results.append({"filename": safe_name, "error": f"No text extractor for '{ext}'"})
            continue

        try:
            extracted_text = await loop.run_in_executor(None, extractor, data)
        except Exception:
            logger.exception("Text extraction failed for file: %s", file.filename)
            results.append({"filename": safe_name, "error": "Failed to extract text"})
            continue

        max_chars = app_settings.ATTACHMENTS_MAX_EXTRACT_CHARS
        truncated = len(extracted_text) > max_chars
        if truncated:
            extracted_text = extracted_text[:max_chars]

        conv_uuid = None
        if conversation_id:
            try:
                conv_uuid = uuid.UUID(conversation_id)
            except ValueError:
                pass
        try:
            upload_record = FileUpload(
                user_id=user_id,
                conversation_id=conv_uuid,
                filename=file.filename,
                extension=ext,
                size_bytes=len(data),
                char_count=len(extracted_text),
                truncated=truncated,
            )
            db.add(upload_record)
            await db.flush()
        except Exception:
            logger.warning("Failed to log upload to DB for %s (table may not exist)", file.filename)
            await db.rollback()

        results.append({
            "filename": safe_name,
            "size": len(data),
            "extension": ext,
            "text": extracted_text,
            "truncated": truncated,
        })

    return {"files": results}


@router.post("/regenerate")
async def regenerate_response(
    body: ChatRequest,
    request: Request,
    user_id: uuid.UUID = Depends(get_current_user_id),
    token=Depends(get_current_user_token),
    db: AsyncSession = Depends(get_db),
):
    """Regenerate the last assistant response in a conversation (streaming)."""
    if not body.conversation_id:
        raise HTTPException(status_code=400, detail="conversation_id is required for regeneration")

    try:
        conv_id = uuid.UUID(body.conversation_id)
    except ValueError:
        raise HTTPException(status_code=400, detail="Invalid conversation_id format")

    await audit_service.log(
        db,
        action="response_regenerated",
        user_id=user_id,
        username=token.username,
        resource_type="conversation",
        resource_id=str(conv_id),
        ip_address=get_client_ip(request),
    )

    return StreamingResponse(
        chat_service.send_message_stream(
            user_id=user_id,
            conversation_id=conv_id,
            content=body.message,
            model=body.model,
        ),
        media_type="application/x-ndjson",
        headers={"Cache-Control": "no-cache", "X-Accel-Buffering": "no"},
    )
